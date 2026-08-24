"""Generates DEFENSE_SLIDES.pptx from the content in DEFENSE_SLIDES.md.

Run: python scripts/generate_defense_pptx.py
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import os

ACCENT = RGBColor(0x0E, 0x63, 0x8C)      # deep cyan/blue — matches app's shield/cyan theme
ACCENT_DARK = RGBColor(0x08, 0x2B, 0x3D)
LIGHT_BG = RGBColor(0xF5, 0xF9, 0xFB)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_title_bar(slide, title, kicker=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if kicker:
        p2 = tf.add_paragraph()
        p2.text = kicker
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xCF, 0xEF, 0xFA)


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def add_bullet_slide(title, bullets, notes=None, kicker=None, numbered=False):
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, title, kicker)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = f"{i + 1}. " if numbered else "•  "
        p.text = prefix + b
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(14)
    add_notes(slide, notes)
    return slide


def add_title_slide(title, subtitle, lines):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_DARK
    bg.line.fill.background()

    box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(0x9F, 0xE0, 0xF2)
    p2.space_before = Pt(10)

    for line in lines:
        p3 = tf.add_paragraph()
        p3.text = line
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0xD8, 0xEE, 0xF6)
        p3.space_before = Pt(6)
    return slide


def add_table_slide(title, headers, rows, extra_bullet=None, notes=None, kicker=None):
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, title, kicker)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.7), Inches(1.5), Inches(11.9), Inches(2.6)
    )
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.size = Pt(16)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(15)
                p.font.color.rgb = TEXT_DARK

    if extra_bullet:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(11.9), Inches(2.4))
        tf = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(extra_bullet):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + b
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(10)
    add_notes(slide, notes)
    return slide


def add_box(slide, left, top, width, height, text, fill=ACCENT, font_color=WHITE, size=13, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = ACCENT_DARK
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_group_label(slide, left, top, width, text):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_DARK
    p.alignment = PP_ALIGN.CENTER


def add_arrow(slide, x1, y1, x2, y2, label=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = GRAY
    conn.line.width = Pt(2)
    conn.line.end_arrowhead = "triangle" if hasattr(conn.line, "end_arrowhead") else None
    if label:
        mid_x, mid_y = (x1 + x2) / 2, (y1 + y2) / 2
        box = slide.shapes.add_textbox(mid_x - Inches(0.9), mid_y - Inches(0.25), Inches(1.8), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
        box.text_frame.word_wrap = True


def add_architecture_diagram_slide():
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, "System Architecture", "Client never talks to third-party APIs directly")

    col_w = Inches(3.6)
    gap = Inches(0.5)
    top0 = Inches(1.9)
    box_h = Inches(0.55)
    box_gap = Inches(0.18)

    col1_x = Inches(0.6)
    col2_x = col1_x + col_w + gap
    col3_x = col2_x + col_w + gap

    add_group_label(slide, col1_x, Inches(1.55), col_w, "📱 Mobile App (Expo)")
    add_group_label(slide, col2_x, Inches(1.55), col_w, "☁️ Supabase")
    add_group_label(slide, col3_x, Inches(1.55), col_w, "🌐 Third-Party APIs")

    device_items = ["UI Screens\n(auth · onboarding · tabs)", "Zustand Stores", "monitoring.ts\n(cycle orchestrator)", "Sensors\n(camera · GPS · audio · accel)"]
    supa_items = ["Auth (email/password)", "Postgres DB\n(Row Level Security)", "Edge Fn: analyse-image / analyse-audio", "Edge Fn: send-sms / send-email / make-call"]
    ext_items = ["OpenAI\nGPT-4o Vision + Whisper", "Twilio\nSMS + Voice", "SendGrid\nEmail", "🆘 Emergency Contact"]

    def add_col(items, x):
        y = top0
        boxes = []
        for item in items:
            b = add_box(slide, x, y, col_w, box_h, item, fill=ACCENT if items is not ext_items else RGBColor(0x2E, 0x8B, 0x57), size=12)
            boxes.append(b)
            y += box_h + box_gap
        return boxes

    device_boxes = add_col(device_items, col1_x)
    supa_boxes = add_col(supa_items, col2_x)
    ext_boxes = add_col(ext_items, col3_x)

    # Recolor emergency contact box distinctly
    ext_boxes[-1].fill.fore_color.rgb = RGBColor(0xB0, 0x30, 0x30)

    mid_y = top0 + Inches(1.2)
    add_arrow(slide, col1_x + col_w, mid_y, col2_x, mid_y, "Edge Function calls\n+ Auth/DB sync")
    add_arrow(slide, col2_x + col_w, mid_y, col3_x, mid_y, "AI analysis\nSMS / Email / Call")

    footer = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.7))
    p = footer.text_frame.paragraphs[0]
    p.text = "Key principle: API keys (OpenAI, Twilio, SendGrid) live only in Supabase Edge Functions — never in the client bundle."
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY
    footer.text_frame.word_wrap = True

    add_notes(slide, "Walk through left to right: device captures sensors, orchestrator decides what to send, Supabase Edge Functions proxy every third-party call so no API key ever ships in the app bundle, and the final alert reaches the emergency contact via Twilio/SendGrid.")
    return slide


# ── Build the deck ────────────────────────────────────────────────────────

add_title_slide(
    "Surveillance AI",
    "AI-Powered Passive Personal Safety Monitoring",
    ["Your Name", "Final Year Project — [Course / Institution]", "Supervisor: [Name]"],
)

add_bullet_slide(
    "The Problem",
    [
        "Personal safety apps today rely on the user actively pressing a button during a crisis",
        "In a real emergency, victims are often unable to act — panicked, restrained, or incapacitated",
        "Existing solutions (panic buttons, manual check-ins) fail exactly when they're needed most",
    ],
    notes="Open with a real scenario: someone walking home alone, or on a first date with a stranger. "
    'Ask: "What if the app could notice danger before you even reach for your phone?"',
)

add_table_slide(
    "The Gap in Existing Solutions",
    ["App type", "Limitation"],
    [
        ["Panic-button apps (Noonlight, bSafe)", "Requires manual trigger"],
        ["Geofencing / location-share apps", "No understanding of what's happening, only where"],
        ["Wearable fall-detection", "Single-signal, no visual/audio context"],
    ],
    extra_bullet=["None of these combine vision + audio + motion into one AI-driven risk assessment that fires automatically"],
)

add_bullet_slide(
    "Solution Overview",
    [
        '"Surveillance AI silently monitors your surroundings using your phone\'s sensors and immediately '
        'alerts your emergency contact by SMS, email, and call if danger is detected."',
        "One-tap start/stop surveillance session",
        "Passive camera snapshots + ambient audio + GPS + accelerometer",
        "AI risk scoring every cycle (Low / Medium / High)",
        "Automatic SMS + email + phone call on High risk",
        "Shake/impact detection bypasses AI for instant alert",
        "Wellness check-in with missed-check-in auto-alert",
    ],
)

add_bullet_slide(
    "Tech Stack",
    [
        "Frontend: Expo, React Native, TypeScript, Expo Router, NativeWind",
        "State: Zustand + AsyncStorage persistence",
        "Backend: Supabase (Postgres + Auth + Row Level Security + Edge Functions)",
        "AI: OpenAI GPT-4o (vision + risk reasoning), Whisper (audio transcription)",
        "Alerts: Twilio (SMS + Voice), SendGrid (Email)",
        "Native sensors: expo-camera, expo-audio, expo-location, expo-sensors, expo-task-manager",
    ],
    notes="Be ready to justify each choice — see the 'why X instead of Y' prep notes.",
)

add_bullet_slide(
    "System Architecture (Overview)",
    [
        "Client never talks to OpenAI/Twilio/SendGrid directly",
        "All third-party API calls proxied through Supabase Edge Functions so API keys never ship inside the app bundle",
        "Client talks to Supabase directly only for Auth and RLS-protected data",
        "See next slide for the full diagram",
    ],
)

add_architecture_diagram_slide()

add_bullet_slide(
    "Onboarding Flow",
    [
        "12-screen guided flow, run once on first launch",
        "Collects: who you're protecting, when you need protection, emergency contact, monitoring interval, shake sensitivity",
        "Requests 4 permissions one at a time with plain-language explanations: Camera → Microphone → Location (Always) → Notifications",
        'Ends with a personalised "plan reveal" screen',
    ],
    notes="Emphasize the UX reasoning: permission requests are proven to have higher acceptance rates when explained individually, in context.",
)

add_bullet_slide(
    "The Monitoring Loop (Core Engine)",
    [
        "Session starts → cycle repeats every 20–30s (configurable)",
        "Camera snapshot + GPS captured in parallel",
        "Event logged immediately with a placeholder — no waiting on AI",
        "Audio recorded concurrently, transcribed via Whisper",
        "Photo sent to GPT-4o Vision for risk analysis",
        "Both AI results combine into one risk score (combineRisks — takes the higher of the two)",
        "Accelerometer shake detection runs independently and continuously — bypasses the AI cycle entirely for instant response",
    ],
    numbered=True,
    notes="This slide is your architecture 'money slide' — know src/lib/monitoring.ts well enough to explain it without notes.",
)

add_bullet_slide(
    "AI Risk Scoring",
    [
        "GPT-4o Vision + Whisper, prompted to return structured JSON: { riskLevel, summary, concerns, confidence }",
        "Prompts explicitly instruct the model to be conservative — only rate High with clear evidence of danger",
        "Medium risk → notifies the user only (human-in-the-loop buffer)",
        "High risk → alerts the emergency contact automatically",
        "Shake events skip AI entirely — the most reliable signal gets the fastest path",
    ],
)

add_bullet_slide(
    "Alert Pipeline",
    [
        "Three channels fire together on High risk: SMS, Email, Phone Call",
        "Phone call reserved for highest-confidence events only (shake + AI High combined, not AI High alone)",
        "All requests routed through Supabase Edge Functions (Twilio, SendGrid)",
        "Deduplication: an alert only fires once per event ID, preventing repeat alerts for the same detected incident",
        "Every alert stores a full record locally and syncs to Supabase",
    ],
)

add_bullet_slide(
    "Security & Privacy",
    [
        "Row Level Security enabled on every Supabase table — auth.uid() scoping means users can only ever access their own data",
        "API keys (OpenAI, Twilio, SendGrid) live only in Edge Functions — never in the client bundle",
        "Emergency contact details and session data stored via Expo Secure Store / RLS-protected tables",
        "Event log auto-clears after 5 days — privacy-by-design retention limit",
        "Open discussion point: recording bystanders raises consent questions under one-party vs two-party consent laws",
    ],
)

add_bullet_slide(
    "Testing & Validation",
    [
        "Manual, scenario-based testing on physical devices (simulators can't test camera, mic, background execution, or push notifications)",
        "Checklist covers: auth flows, onboarding, permission-denied degraded mode, session lifecycle, background execution, shake detection, alert firing",
        "AI accuracy was evaluated qualitatively against staged scenarios, not a formal labeled dataset (explicitly named as a limitation)",
    ],
)

add_bullet_slide(
    "Limitations & Future Work",
    [
        "No automated test suite / no formal precision-recall evaluation of AI risk scoring",
        "Battery and data usage under continuous monitoring not yet profiled",
        "No offline fallback if network or AI provider is unreachable mid-emergency",
        "Single emergency contact only — no escalation chain",
        "Consent/legal questions around recording bystanders unresolved",
        "Future: on-device fallback model for offline risk scoring, multiple contacts with escalation tiers, wearable integration",
    ],
)

add_bullet_slide(
    "Live Demo",
    [
        "Start a surveillance session → show event appearing in the Log tab",
        "Trigger a simulated shake event → show instant High-risk alert path",
        "Show Settings: emergency contact, monitoring interval, shake sensitivity",
        "Have a backup screen recording ready in case live demo fails",
    ],
)

add_bullet_slide(
    "Conclusion & Q&A",
    [
        "Recap the one-liner and the core engineering contribution: combining multimodal AI signals into one automatic, passive alert pipeline",
        "Thank the panel, open the floor",
    ],
)

out_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "DEFENSE_SLIDES.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
